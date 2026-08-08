import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_impressive_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Tokens
    BG_MAIN = RGBColor(11, 15, 25)       # Rich Deep Navy #0b0f19
    CARD_BG = RGBColor(22, 30, 49)       # Glass Navy #161e31
    CARD_SUB = RGBColor(16, 23, 39)      # Darker Subcard #101727
    BORDER_COL = RGBColor(40, 56, 89)    # Border Slate
    
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(156, 163, 175) # Gray 400
    TEXT_DIM = RGBColor(203, 213, 225)   # Slate 300
    
    CYAN = RGBColor(56, 189, 248)        # Sky 400
    EMERALD = RGBColor(52, 211, 153)     # Emerald 400
    AMBER = RGBColor(251, 191, 36)       # Amber 400
    ROSE = RGBColor(251, 113, 133)       # Rose 400
    PURPLE = RGBColor(192, 132, 252)     # Purple 400
    INDIGO = RGBColor(129, 140, 248)     # Indigo 400

    def set_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_MAIN
        bg.line.fill.background()
        return bg

    def add_card(slide, l, t, w, h, bg=CARD_BG, border=BORDER_COL, border_w=1):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = bg
        if border:
            card.line.color.rgb = border
            card.line.width = Pt(border_w)
        else:
            card.line.fill.background()
        return card

    def add_badge(slide, l, t, text, bg_col, text_col):
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, Inches(len(text)*0.11 + 0.5), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = bg_col
        badge.line.color.rgb = text_col
        badge.line.width = Pt(1)
        tf = badge.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = text
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = text_col
        return badge

    def add_header(slide, tag_text, title, subtitle):
        add_badge(slide, Inches(0.6), Inches(0.4), tag_text, RGBColor(20, 35, 60), CYAN)
        
        # Title + Subtitle box
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(12.133), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 1: Title & Hero
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1)

    add_card(s1, Inches(0.6), Inches(0.6), Inches(12.133), Inches(6.3), CARD_BG, CYAN, 2)
    add_badge(s1, Inches(1.0), Inches(1.0), "🚀 NEXT-GENERATION RECRUITMENT AI", RGBColor(16, 40, 60), EMERALD)

    tb = s1.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(11.333), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Resume Screening Agent"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "Collaborative Multi-Agent Resume Screening & Career Intelligence Platform"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = CYAN

    p3 = tf.add_paragraph()
    p3.text = "Solving the resume 'black hole' with 4 transparent, cooperating AI agents. Designed for job seekers, learners, and high-volume recruiters."
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_DIM

    metrics = [
        ("⏱️ 95% Time Saved", "Screens 50+ candidates in <30s instead of 2 weeks of manual reading.", CYAN, Inches(1.0)),
        ("🤝 4 AI Agents", "Intake Reader, Reviewer, Fact-Checking QA & Interview Coach.", EMERALD, Inches(3.85)),
        ("🛡️ Zero Hallucinations", "Math vector verification (TF-IDF) stops fake scores and bias.", AMBER, Inches(6.7)),
        ("🎯 Free Interview Prep", "Auto-generates tailored technical and behavioral practice questions.", PURPLE, Inches(9.55)),
    ]
    for title, desc, col, l in metrics:
        add_card(s1, l, Inches(4.3), Inches(2.75), Inches(1.9), CARD_SUB, col, 1.5)
        tb = s1.shapes.add_textbox(l + Inches(0.12), Inches(4.4), Inches(2.5), Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = "\n" + desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_WHITE

    tb_f = s1.shapes.add_textbox(Inches(1.0), Inches(6.4), Inches(11.333), Inches(0.4))
    tf_f = tb_f.text_frame
    p = tf_f.paragraphs[0]
    p.text = "Project: Resume Screening Agent | Stack: Python 3.13, Google Gemini 3.5 Flash, Scikit-Learn, Streamlit"
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 2: The Real-Life Problem & Industry Gaps
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_header(s2, "THE REAL-WORLD PROBLEM", "The Broken 'Black Hole' of Modern Job Applications", 
               "Why 75% of qualified resumes get dropped before a human ever sees them.")

    prob_cards = [
        ("👤 Job Seekers & Learners", ROSE, [
            "• Spend 10+ hours tailoring resumes with no feedback.",
            "• Discarded by black-box algorithms with zero explanation.",
            "• No way to know which missing skills caused rejection.",
            "• Leaves candidates frustrated, confused, and stuck."
        ], Inches(0.6)),
        ("🤖 Traditional Dumb ATS Filters", AMBER, [
            "• Rigid keyword matching (like a blind 'Ctrl + F' tool).",
            "• Example: Rejects 'Client Help' if JD says 'Customer Care'.",
            "• Easily tricked by 'white-text' keyword stuffing.",
            "• Fails to understand real-world projects or career depth."
        ], Inches(4.7)),
        ("💼 Overwhelmed Recruiters", CYAN, [
            "• 500+ applications per open position within 48 hours.",
            "• Average 6 seconds spent per resume manual scan.",
            "• Single-prompt AI tools hallucinate and give fake 99% scores.",
            "• Takes 2 to 3 weeks just to build a reliable shortlist."
        ], Inches(8.8)),
    ]
    for title, col, points, l in prob_cards:
        add_card(s2, l, Inches(1.8), Inches(3.93), Inches(4.5), CARD_BG, col, 1.5)
        tb = s2.shapes.add_textbox(l + Inches(0.18), Inches(1.95), Inches(3.55), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = col

        for pt in points:
            p_pt = tf.add_paragraph()
            p_pt.text = "\n" + pt
            p_pt.font.size = Pt(12)
            p_pt.font.color.rgb = TEXT_WHITE

    add_card(s2, Inches(0.6), Inches(6.45), Inches(12.133), Inches(0.65), RGBColor(35, 20, 30), ROSE, 1)
    tb_r = s2.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.55))
    tf_r = tb_r.text_frame
    p = tf_r.paragraphs[0]
    p.text = "⚡ Core Flaw: Old tools test for EXACT KEYWORDS rather than TRUE COMPETENCE & VALUE."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ROSE

    # =========================================================================
    # SLIDE 3: ATS & Scores Explained in Plain English (Analogies)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_header(s3, "PLAIN ENGLISH EXPLANATION", "Demystifying ATS & Resume Scores for Non-Tech Learners",
               "Understanding automated screening through intuitive real-world comparisons.")

    add_card(s3, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.2), CARD_BG, EMERALD, 1.5)
    tb_l = s3.shapes.add_textbox(Inches(0.8), Inches(1.95), Inches(5.5), Inches(4.9))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "🍰 The Master Chef & Recipe Analogy"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    chef_items = [
        "1. Job Description = The Chef's Signature Recipe\n   (Needs: 3 cups flour, chocolate, eggs, 5 yrs baking mastery)",
        "2. Your Resume = Your Basket of Ingredients\n   (What skills, practical projects, and work history you bring)",
        "3. The ATS System = The Kitchen Inspector\n   (Checks whether your basket contains what the recipe needs)",
        "4. The 0-100% Score = The Match Percentage\n   (How complete and fresh your ingredients are for this dish)"
    ]
    for item in chef_items:
        p_i = tf_l.add_paragraph()
        p_i.text = "\n" + item
        p_i.font.size = Pt(12)
        p_i.font.color.rgb = TEXT_WHITE

    add_card(s3, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.2), CARD_BG, CYAN, 1.5)
    tb_r = s3.shapes.add_textbox(Inches(7.0), Inches(1.95), Inches(5.5), Inches(4.9))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "⚖️ Old Keyword Robots vs. Resume Screening Agent"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN

    vs_items = [
        "❌ Old Dumb Robots (Literal Word Match):\n• If JD asks for 'FastAPI' and you wrote 'REST APIs with Python', it marks 0%!\n• Easily confused by resume formatting or synonyms.\n• Rejects honest learners with real potential.",
        "✅ Resume Screening Agent (Context & Meaning Aware):\n• Understands semantic meaning: knows 'REST APIs with Python' relates directly to backend engineering.\n• Multi-agent cross-verification stops false rejections."
    ]
    for item in vs_items:
        p_v = tf_r.add_paragraph()
        p_v.text = "\n" + item
        p_v.font.size = Pt(12)
        p_v.font.color.rgb = TEXT_WHITE if not "❌" in item else ROSE

    # =========================================================================
    # SLIDE 4: Architecture — The 4 Collaborative AI Agents
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_header(s4, "COLLABORATIVE AI PIPELINE", "The 4-Agent Hiring Committee: How They Work Together",
               "Why 4 specialized agents working in sequence are 10x more reliable than a single AI prompt.")

    agents = [
        ("Agent 1: The Reader 🔍", "Resume Parser Agent", CYAN, [
            "• Extracts clean data from PDF, Word (.docx), & TXT.",
            "• Extracts: Work History, Timeline, Skills, Education.",
            "• Normalizes messy tables & columns into clean JSON."
        ], Inches(0.6)),
        ("Agent 2: The Grader ⚖️", "Candidate Review Agent", EMERALD, [
            "• Scores 4 core pillars (0-100 each):",
            "  1. Technical Skills  2. Experience Depth",
            "  3. Education Fit     4. Role Match",
            "• Identifies concrete strengths & skill gaps."
        ], Inches(3.7)),
        ("Agent 3: The Auditor 🛡️", "Quality Assurance Agent", AMBER, [
            "• Fact-checks Agent 2 against parsed timeline.",
            "• Incorporates ML vector similarity (TF-IDF).",
            "• Penalizes keyword-stuffing and fixes score inflation.",
            "• Writes transparent audit justifications."
        ], Inches(6.8)),
        ("Agent 4: The Coach 👑", "Ranker & Interviewer", PURPLE, [
            "• Compiles full leaderboard with final ranks.",
            "• Generates 3-4 custom interview questions tailored to each candidate's specific gaps.",
            "• Produces executive batch summary for HR."
        ], Inches(9.9)),
    ]
    for title, sub, col, points, l in agents:
        add_card(s4, l, Inches(1.8), Inches(2.85), Inches(5.2), CARD_BG, col, 1.5)
        tb = s4.shapes.add_textbox(l + Inches(0.12), Inches(1.95), Inches(2.6), Inches(4.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col

        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = TEXT_MUTED

        for pt in points:
            p_pt = tf.add_paragraph()
            p_pt.text = "\n" + pt
            p_pt.font.size = Pt(11.5)
            p_pt.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 5: Step-by-Step Workflow & Data Flow
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_header(s5, "END-TO-END WORKFLOW", "From Raw Document Upload to Actionable Report in <30s",
               "A fast, deterministic 5-stage pipeline designed for simplicity and auditability.")

    flow_steps = [
        ("1. Ingestion 📥", "Recruiter drops 10-100 resumes + pastes Job Description.", CYAN, Inches(0.6), Inches(1.8)),
        ("2. Text Extraction 📄", "PyPDF and python-docx parse clean text, ignoring bad formatting.", INDIGO, Inches(6.8), Inches(1.8)),
        ("3. Math Baseline 📊", "Scikit-Learn calculates TF-IDF & Cosine Similarity vector scores.", EMERALD, Inches(0.6), Inches(3.6)),
        ("4. Multi-Agent Audit 🤖", "Parser ➔ Evaluator ➔ QA Auditor ➔ Ranker execute cooperatively.", AMBER, Inches(6.8), Inches(3.6)),
        ("5. Rich Output Dashboard 🎯", "Export CSV/JSON, view candidate leaderboard, and practice generated interview questions.", PURPLE, Inches(0.6), Inches(5.4)),
    ]
    for title, desc, col, l, t in flow_steps[:4]:
        add_card(s5, l, t, Inches(5.9), Inches(1.6), CARD_BG, col, 1.5)
        tb = s5.shapes.add_textbox(l + Inches(0.18), t + Inches(0.1), Inches(5.5), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col
        p_d = tf.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_WHITE

    add_card(s5, flow_steps[4][3], flow_steps[4][4], Inches(12.133), Inches(1.55), CARD_BG, flow_steps[4][2], 1.5)
    tb_5 = s5.shapes.add_textbox(flow_steps[4][3] + Inches(0.18), flow_steps[4][4] + Inches(0.1), Inches(11.7), Inches(1.3))
    tf_5 = tb_5.text_frame
    tf_5.word_wrap = True
    p = tf_5.paragraphs[0]
    p.text = flow_steps[4][0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = flow_steps[4][2]
    p_d = tf_5.add_paragraph()
    p_d.text = "\n" + flow_steps[4][1]
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 6: Real Candidate Case Study (Alice, Charlie, Diana)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_header(s6, "TESTED BENCHMARK RESULTS", "Case Study: Screening Candidates for a Senior Full-Stack Role",
               "Target Job: Senior Full-Stack Engineer (Python, React, FastAPI, Docker, AWS)")

    candidates = [
        ("🏆 Alice Smith", "Senior Full-Stack Lead", "92 / 100", EMERALD, [
            "• Profile: 7+ yrs Python, React, FastAPI, AWS, Docker.",
            "• QA Audit: Verified real architecture leadership.",
            "• Decision: SHORTLIST FOR FINAL ROUND.",
            "• Custom Interview Q: 'How did you handle Redis caching & database contention at peak load?'"
        ], Inches(0.6)),
        ("⚠️ Charlie Brown", "Junior Frontend Learner", "45 / 100", AMBER, [
            "• Profile: 1 yr HTML/CSS, basic React projects.",
            "• QA Audit: Missing Python backend & cloud infra.",
            "• Decision: NOT READY FOR SENIOR ROLE.",
            "• Actionable Advice: Build FastAPI backends & Dockerize apps before applying to senior roles."
        ], Inches(4.7)),
        ("❌ Diana Prince", "Project Manager (Non-Tech)", "10 / 100", ROSE, [
            "• Profile: PMP, Agile Scrum, Jira, budgeting.",
            "• QA Audit: Zero hands-on coding or API work.",
            "• Decision: REJECT FOR CODING OPENING.",
            "• Actionable Advice: Re-route resume to Technical Project Management positions."
        ], Inches(8.8)),
    ]
    for name, role, score, col, points, l in candidates:
        add_card(s6, l, Inches(1.8), Inches(3.93), Inches(5.2), CARD_BG, col, 1.5)
        tb = s6.shapes.add_textbox(l + Inches(0.18), Inches(1.95), Inches(3.55), Inches(4.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"{name} ({score})"
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = col

        p_sub = tf.add_paragraph()
        p_sub.text = role
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = TEXT_MUTED

        for pt in points:
            p_pt = tf.add_paragraph()
            p_pt.text = "\n" + pt
            p_pt.font.size = Pt(11.5)
            p_pt.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 7: Empowering Job Seekers & Learners
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_header(s7, "FOR CANDIDATES & STUDENTS", "How Resume Screening Agent Accelerates Your Career Growth",
               "Turning rejection into a personalized roadmap for interview preparation.")

    seeker_cards = [
        ("1. 100% Transparent Scoring 🎯", CYAN, 
         "No hidden black-box magic. See exact category marks for Technical Skills, Experience, Education, and Role Alignment so you know where you stand."),
        ("2. Actionable Gap Analysis 💡", EMERALD, 
         "Instead of generic 'we went with another candidate' emails, see the exact libraries, tools, or concepts you should learn next to land the role."),
        ("3. Free Tailored Interview Prep 🎙️", PURPLE, 
         "The AI generates 3-4 custom technical and behavioral questions targeted at your exact background so you can practice before live interviews."),
        ("4. Fair & Unbiased Fact-Checking 🛡️", AMBER, 
         "No penalty for layout formatting, graphics, or fonts. Our multi-agent auditor rewards genuine skills and real experience."),
    ]
    for idx, (title, col, desc) in enumerate(seeker_cards):
        col_idx = idx % 2
        row_idx = idx // 2
        l = Inches(0.6 + col_idx * 6.2)
        t = Inches(1.8 + row_idx * 2.65)
        add_card(s7, l, t, Inches(5.9), Inches(2.45), CARD_BG, col, 1.5)
        tb = s7.shapes.add_textbox(l + Inches(0.18), t + Inches(0.15), Inches(5.5), Inches(2.15))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col
        p_d = tf.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.size = Pt(12.5)
        p_d.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 8: Benefits for Recruiters & Enterprises
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8)
    add_header(s8, "FOR RECRUITERS & ENTERPRISES", "Enterprise-Grade Speed, Accuracy & Auditability",
               "Helping talent acquisition teams screen 50+ candidates in under 30 seconds.")

    rec_grid = [
        ("⏱️ 95% Time Reduction", "Slash candidate screening time from 2 weeks down to 30 seconds.", CYAN),
        ("🛡️ Anti-Hallucination QA", "Cross-checks LLM scores against TF-IDF mathematical cosine similarity.", AMBER),
        ("📊 Multi-Format Exports", "One-click export to CSV (Excel), JSON, and Executive Markdown summaries.", EMERALD),
        ("🎯 Automated Interview Kits", "Generates targeted interview question kits for hiring managers automatically.", PURPLE),
        ("🌐 Multi-Environment", "Runs on Local Web UI, Command Line Interface (CLI), or Docker containers.", INDIGO),
        ("💰 Near-Zero Infrastructure Cost", "Built with lightweight Python & ultra-fast Google Gemini 3.5 Flash.", ROSE),
    ]
    for idx, (title, desc, col) in enumerate(rec_grid):
        col_idx = idx % 3
        row_idx = idx // 3
        l = Inches(0.6 + col_idx * 4.1)
        t = Inches(1.8 + row_idx * 2.65)
        add_card(s8, l, t, Inches(3.93), Inches(2.45), CARD_BG, col, 1.5)
        tb = s8.shapes.add_textbox(l + Inches(0.15), t + Inches(0.15), Inches(3.6), Inches(2.15))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col
        p_d = tf.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 9: User-Friendly Dashboard Interface
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9)
    add_header(s9, "STREAMLIT DASHBOARD", "Intuitive Dark-Themed UI: Zero Technical Setup Needed",
               "Designed for ease-of-use with instant drag-and-drop file processing.")

    add_card(s9, Inches(0.6), Inches(1.8), Inches(5.9), Inches(5.2), CARD_BG, CYAN, 1.5)
    tb_l = s9.shapes.add_textbox(Inches(0.8), Inches(1.95), Inches(5.5), Inches(4.9))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "🖥️ Simple 3-Step User Experience"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN

    ui_steps = [
        "1. Paste Job Description:\n   Pre-loaded with sample roles or custom paste any job opening.",
        "2. Drop Candidate Resumes:\n   Supports batch uploads for PDF, Word (.docx), and plain text (.txt).",
        "3. Live Real-Time Execution:\n   Visual progress bars and live agent audit logs display every step as it happens.",
        "4. Explore & Export Leaderboard:\n   View ranked scores, click candidate tabs, and download results to Excel."
    ]
    for step in ui_steps:
        p_s = tf_l.add_paragraph()
        p_s.text = "\n" + step
        p_s.font.size = Pt(12)
        p_s.font.color.rgb = TEXT_WHITE

    add_card(s9, Inches(6.8), Inches(1.8), Inches(5.9), Inches(5.2), CARD_BG, EMERALD, 1.5)
    tb_r = s9.shapes.add_textbox(Inches(7.0), Inches(1.95), Inches(5.5), Inches(4.9))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "📊 Interactive Analytics Panels"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = EMERALD

    panels = [
        "• 📈 Score Visualizer: Horizontal score comparison bars for fast decision making.",
        "• 🎯 Strengths & Gaps View: Color-coded bullet points showing missing qualifications.",
        "• 🎙️ Tailored Question Drawer: Candidate-specific technical and behavioral questions.",
        "• 🛡️ QA Audit Logs: Complete visibility into why scores were adjusted."
    ]
    for p_item in panels:
        p_p = tf_r.add_paragraph()
        p_p.text = "\n" + p_item
        p_p.font.size = Pt(12)
        p_p.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 10: The Engineering & Tech Stack (Explained Simply)
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10)
    add_header(s10, "TECHNOLOGY STACK", "Modern, Lightweight & Scalable Tech Stack",
               "Built on state-of-the-art AI, machine learning, and Python standards.")

    stack_cards = [
        ("🧠 Google Gemini AI", "gemini-3.5-flash", CYAN, "Ultra-fast LLM that reads resumes with human-level context and outputs structured Pydantic JSON payloads."),
        ("📊 Scikit-Learn (ML)", "TF-IDF + Cosine Sim", EMERALD, "Mathematical vectorization of text that measures vocabulary overlap to anchor the AI's grading."),
        ("💻 Streamlit Framework", "Interactive Web UI", AMBER, "Lightweight Python web dashboard that provides reactive, real-time feedback with zero JavaScript bloat."),
        ("🛡️ Pydantic & Pytest", "Data Validation & CI", PURPLE, "Enforces strict JSON schema integrity with automated unit testing suites in GitHub Actions."),
    ]
    for idx, (title, sub, col, desc) in enumerate(stack_cards):
        col_idx = idx % 2
        row_idx = idx // 2
        l = Inches(0.6 + col_idx * 6.2)
        t = Inches(1.8 + row_idx * 2.65)
        add_card(s10, l, t, Inches(5.9), Inches(2.45), CARD_BG, col, 1.5)
        tb = s10.shapes.add_textbox(l + Inches(0.18), t + Inches(0.15), Inches(5.5), Inches(2.15))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col

        p_s = tf.add_paragraph()
        p_s.text = sub
        p_s.font.size = Pt(11)
        p_s.font.color.rgb = TEXT_MUTED

        p_d = tf.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 11: Summary & Key Takeaways
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_bg(s11)
    add_header(s11, "EXECUTIVE SUMMARY", "Key Takeaways: Why Resume Screening Agent Matters",
               "Transforming resume screening from a mystery into an empowering, transparent tool.")

    sum_points = [
        ("✨ 100% Transparent & Fair", "Replaces blind keyword counting with contextual semantic understanding and multi-agent fact-checking.", CYAN),
        ("🤝 4-Agent Cooperation", "Parser, Reviewer, QA Auditor, and Coach ensure no single AI hallucination affects a candidate's future.", EMERALD),
        ("🚀 Career Growth for Learners", "Job seekers get exact skill gap diagnostics and tailored interview questions to practice.", PURPLE),
        ("⚡ 95% Time Saved for HR", "Recruiters screen large batches in under 30 seconds with complete exportable audit records.", AMBER),
    ]
    for idx, (title, desc, col) in enumerate(sum_points):
        col_idx = idx % 2
        row_idx = idx // 2
        l = Inches(0.6 + col_idx * 6.2)
        t = Inches(1.8 + row_idx * 2.65)
        add_card(s11, l, t, Inches(5.9), Inches(2.45), CARD_BG, col, 1.5)
        tb = s11.shapes.add_textbox(l + Inches(0.18), t + Inches(0.15), Inches(5.5), Inches(2.15))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col
        p_d = tf.add_paragraph()
        p_d.text = "\n" + desc
        p_d.font.size = Pt(12.5)
        p_d.font.color.rgb = TEXT_WHITE

    # =========================================================================
    # SLIDE 12: Q&A, Live Demo & Contact
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_bg(s12)
    add_card(s12, Inches(0.6), Inches(0.6), Inches(12.133), Inches(6.3), CARD_BG, EMERALD, 2)

    add_badge(s12, Inches(1.0), Inches(1.0), "THANK YOU & DISCUSSION", RGBColor(16, 40, 30), EMERALD)

    tb = s12.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Thank You! 🙏"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    p2 = tf.add_paragraph()
    p2.text = "Open for Questions, Live Demo & Discussion"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = CYAN

    links = [
        "\n🌐 Live Streamlit Dashboard: http://localhost:8501",
        "📂 GitHub Repository: Snehan0911/Resume-Screening-Agent",
        "🐳 Docker Deployment: docker run -p 8501:8501 resume-screening-agent",
        "\n💬 'Empowering candidates with clear feedback while helping companies hire the right talent faster.'"
    ]
    for lk in links:
        p_l = tf.add_paragraph()
        p_l.text = lk
        p_l.font.size = Pt(13.5)
        p_l.font.color.rgb = TEXT_WHITE if not "💬" in lk else TEXT_MUTED

    output_path = os.path.join(os.getcwd(), "Resume_Screening_Agent_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation successfully saved to: {output_path}")

if __name__ == "__main__":
    build_impressive_deck()
